import io
import logging
from typing import NamedTuple

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

logger = logging.getLogger(__name__)

_MAX_DOCUMENT_BYTES = 10 * 1024 * 1024

_ALLOWED_EXTENSIONS = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc": "application/msword",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt": "application/vnd.ms-powerpoint",
    "txt": "text/plain",
    "text": "text/plain",
    "md": "text/markdown",
    "markdown": "text/markdown",
    "csv": "text/csv",
}


class ExtractedDocument(NamedTuple):
    text: str
    stored_as: str


def _extension(filename: str | None) -> str:
    if not filename:
        return ""
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def extract_document_text(filename: str | None, content: bytes) -> str:
    if len(content) > _MAX_DOCUMENT_BYTES:
        raise ValueError("Document exceeds the 10 MB size limit")

    ext = _extension(filename)
    if ext not in _ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported document type '.{ext or 'unknown'}'. "
            f"Allowed: {', '.join(sorted(_ALLOWED_EXTENSIONS))}"
        )

    try:
        if ext == "pdf":
            reader = PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext in ("docx", "doc"):
            document = DocxDocument(io.BytesIO(content))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        elif ext in ("pptx", "ppt"):
            presentation = Presentation(io.BytesIO(content))
            parts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            text = "\n".join(parts)
        else:
            text = content.decode("utf-8", errors="replace")
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("Failed to parse document: %s", exc)
        raise ValueError(f"Could not read the document content: {exc}")

    text = text.strip()
    if not text:
        raise ValueError("No readable text could be extracted from the document")
    return text
